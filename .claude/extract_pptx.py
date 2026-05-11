"""Extract text from PPTX files for IR synthesis (UTF-8 output)."""
import sys
import io
from pptx import Presentation

# Force UTF-8 stdout
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_pptx(path: str):
    print(f"\n{'='*80}")
    print(f"FILE: {path}")
    print(f"{'='*80}\n")
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n## Slide {i}\n")
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(" |"):
                            texts.append(f"[TABLE] {row_text}")
            if texts:
                print("\n".join(texts))
            else:
                print("(no text)")
    except Exception as e:
        print(f"ERROR: {e}")

if __name__ == "__main__":
    for path in sys.argv[1:]:
        extract_pptx(path)
